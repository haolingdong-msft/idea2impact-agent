import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


RUN_DIR = Path(sys.argv[1])
OUTPUT = Path(sys.argv[2])


def color(value: str) -> RGBColor:
    return RGBColor.from_string(value.lstrip("#"))


def add_page(presentation: Presentation, page_dir: Path) -> int:
    manifest = json.loads((page_dir / "manifest.json").read_text(encoding="utf-8"))
    slide_spec = manifest["slide"]
    content = manifest["content_box"]
    source = manifest["source"]
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color(
        slide_spec.get("background", "#FFFFFF")
    )

    def box(box_px):
        x, y, width, height = box_px
        return (
            Inches(content["left"] + x / source["width_px"] * content["width"]),
            Inches(content["top"] + y / source["height_px"] * content["height"]),
            Inches(width / source["width_px"] * content["width"]),
            Inches(height / source["height_px"] * content["height"]),
        )

    objects = []
    for index, item in enumerate(manifest.get("shapes", [])):
        objects.append((item.get("z_index", 100), index, "shape", item))
    for index, item in enumerate(manifest.get("images", [])):
        objects.append((item.get("z_index", 200), index, "image", item))
    for index, item in enumerate(manifest.get("text_boxes", [])):
        objects.append((item.get("z_index", 300), index, "text", item))

    for _, _, kind, item in sorted(objects):
        if kind == "shape":
            if item["type"] == "line":
                x1, y1, x2, y2 = item["points_px"]
                left, top, _, _ = box([x1, y1, 0, 0])
                right, bottom, _, _ = box([x2, y2, 0, 0])
                shape = slide.shapes.add_connector(
                    MSO_CONNECTOR.STRAIGHT, left, top, right, bottom
                )
                shape.line.color.rgb = color(item.get("stroke", "#000000"))
                shape.line.width = Pt(item.get("stroke_width", 1))
                continue
            shape_type = {
                "rect": MSO_SHAPE.RECTANGLE,
                "roundRect": MSO_SHAPE.ROUNDED_RECTANGLE,
                "ellipse": MSO_SHAPE.OVAL,
            }.get(item["type"], MSO_SHAPE.RECTANGLE)
            shape = slide.shapes.add_shape(shape_type, *box(item["box_px"]))
            fill = item.get("fill")
            if fill and fill != "none":
                shape.fill.solid()
                shape.fill.fore_color.rgb = color(fill)
            else:
                shape.fill.background()
            stroke = item.get("stroke")
            if stroke and stroke != "none":
                shape.line.color.rgb = color(stroke)
                shape.line.width = Pt(item.get("stroke_width", 1))
            else:
                shape.line.fill.background()
            continue

        if kind == "image":
            slide.shapes.add_picture(
                str(page_dir / item["path"]), *box(item["box_px"])
            )
            continue

        text_box = slide.shapes.add_textbox(*box(item["box_px"]))
        frame = text_box.text_frame
        frame.clear()
        frame.margin_left = 0
        frame.margin_right = 0
        frame.margin_top = 0
        frame.margin_bottom = 0
        frame.word_wrap = bool(item.get("fit_text", True))
        frame.vertical_anchor = {
            "middle": MSO_ANCHOR.MIDDLE,
            "bottom": MSO_ANCHOR.BOTTOM,
        }.get(item.get("valign"), MSO_ANCHOR.TOP)
        for line_index, line in enumerate(str(item.get("text", "")).split("\n")):
            paragraph = (
                frame.paragraphs[0] if line_index == 0 else frame.add_paragraph()
            )
            paragraph.alignment = {
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
            }.get(item.get("align"), PP_ALIGN.LEFT)
            paragraph.space_before = Pt(0)
            paragraph.space_after = Pt(0)
            run = paragraph.add_run()
            run.text = line
            run.font.name = item.get("font", "Arial")
            run.font.size = Pt(item.get("font_size", 10))
            run.font.bold = bool(item.get("bold", False))
            run.font.italic = bool(item.get("italic", False))
            run.font.color.rgb = color(item.get("color", "#111111"))
    return len(slide.shapes)


deck = json.loads((RUN_DIR / "deck_manifest.json").read_text(encoding="utf-8"))
page_dirs = [
    RUN_DIR / page["page_dir"] if isinstance(page, dict) else RUN_DIR / page
    for page in deck["pages"]
]
first = json.loads((page_dirs[0] / "manifest.json").read_text(encoding="utf-8"))
presentation = Presentation()
presentation.slide_width = Inches(first["slide"]["width"])
presentation.slide_height = Inches(first["slide"]["height"])
counts = [add_page(presentation, page_dir) for page_dir in page_dirs]
presentation.save(OUTPUT)
print(json.dumps({"output": str(OUTPUT), "slides": len(counts), "shapes": counts}))
